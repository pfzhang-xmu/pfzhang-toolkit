#!/usr/bin/env python3
"""Build a single-slide PPTX from a JSON specification.

Usage:
    python build_slide.py --spec PATH --output PATH [--template PATH]

JSON Specification Schema:
{
  "slide": {
    "width_inches": 10.0,
    "height_inches": 5.625,
    "background": {"type": "solid", "color": "#FFFFFF"}
  },
  "elements": [
    {"type": "text_box", "x": 1.0, "y": 0.5, "w": 5.0, "h": 1.0,
     "runs": [{"text": "...", "font": "PingFang SC", "size_pt": 36,
               "bold": true, "italic": false, "color": "#2C3E50"}],
     "alignment": "left", "valign": "top", "margin_pt": 0,
     "line_spacing_pt": null, "layer": 0},

    {"type": "image", "path": "/tmp/crops/logo.png", "x": 0.5, "y": 0.3,
     "w": 1.2, "h": 0.8, "layer": 0},

    {"type": "shape", "subtype": "rectangle", "x": 0, "y": 0, "w": 10, "h": 5.625,
     "fill": "#FFFFFF", "border": null, "layer": -1},

    {"type": "shape", "subtype": "rounded_rectangle", "x": 0.5, "y": 0.3,
     "w": 3.0, "h": 0.8, "fill": "#3498DB",
     "border": {"color": "#2980B9", "width_pt": 1.0},
     "corner_radius_inches": 0.05,
     "shadow": {"blur_pt": 4, "offset_x_pt": 2, "offset_y_pt": 2, "color": "#00000040"},
     "layer": 2},

    {"type": "line", "x1": 0.5, "y1": 1.5, "x2": 9.5, "y2": 1.5,
     "color": "#E0E0E0", "width_pt": 1.5, "layer": 1},

    {"type": "simple_table", "x": 0.5, "y": 2.0, "w": 9.0, "h": 3.0,
     "rows": 5, "cols": 4,
     "cells": [{"row": 0, "col": 0, "text": "Revenue", "font": "PingFang SC",
                "size_pt": 11, "bold": true, "fill": "#2C3E50", "color": "#FFFFFF"}],
     "border_color": "#BDC3C7", "border_width_pt": 0.5,
     "layer": 3}
  ]
}
"""

import argparse
import json
import sys
from copy import deepcopy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn, nsmap
except ImportError:
    print("ERROR: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def hex_to_rgb(hex_color):
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    elif len(hex_color) == 8:
        # 8-char hex with alpha: ignore alpha for fill
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    raise ValueError(f"Invalid hex color: {hex_color}")


def px_to_inches(px, dpi=150):
    """Convert pixels to inches at given DPI."""
    return px / dpi


def inches_to_emu(inches):
    """Convert inches to EMU (English Metric Units)."""
    return int(inches * 914400)


def set_slide_background(slide, bg_spec):
    """Set slide background from spec."""
    if bg_spec.get("type") == "solid":
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_spec["color"])
    elif bg_spec.get("type") == "gradient":
        # Gradient background via XML injection
        bg = slide.background
        fill = bg.fill
        fill.solid()  # Fallback
        # For gradient, we'd inject <p:bg> XML - defer to a later pass if needed


def add_text_box(slide, spec):
    """Add a text box to the slide."""
    x = Inches(spec["x"])
    y = Inches(spec["y"])
    w = Inches(spec["w"])
    h = Inches(spec["h"])

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Margins
    margin = Pt(spec.get("margin_pt", 2))
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Pt(spec.get("margin_top_pt", spec.get("margin_pt", 2)))
    tf.margin_bottom = Pt(spec.get("margin_bottom_pt", spec.get("margin_pt", 2)))

    # Process runs
    runs = spec.get("runs", [])
    if not runs and spec.get("text"):
        runs = [{"text": spec["text"]}]

    for i, run_spec in enumerate(runs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }.get(spec.get("alignment", "left"), PP_ALIGN.LEFT)

        # Line spacing
        if spec.get("line_spacing_pt"):
            p.line_spacing = Pt(spec["line_spacing_pt"])

        run = p.add_run()
        run.text = run_spec.get("text", "")

        # Font
        font = run_spec.get("font", "PingFang SC")
        run.font.name = font

        # Set East Asian font
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.makeelement(qn('a:ea'), {})
            ea.set('typeface', font)
            rPr.append(ea)
        except Exception:
            pass

        run.font.size = Pt(run_spec.get("size_pt", 14))
        run.font.bold = run_spec.get("bold", False)
        run.font.italic = run_spec.get("italic", False)

        if "color" in run_spec:
            run.font.color.rgb = hex_to_rgb(run_spec["color"])

    # Vertical alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)

    return txBox


def add_image(slide, spec):
    """Add an image to the slide."""
    path = spec["path"]
    if not Path(path).exists():
        print(f"WARNING: Image not found: {path}", file=sys.stderr)
        # Create a placeholder rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(spec["x"]), Inches(spec["y"]),
            Inches(spec["w"]), Inches(spec["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        return shape

    try:
        pic = slide.shapes.add_picture(
            path, Inches(spec["x"]), Inches(spec["y"]),
            Inches(spec["w"]), Inches(spec["h"])
        )
        return pic
    except Exception as e:
        print(f"WARNING: Failed to add image {path}: {e}", file=sys.stderr)
        return None


def add_shape(slide, spec):
    """Add a shape (rectangle, rounded rectangle, oval, etc.)."""
    subtype = spec.get("subtype", "rectangle")
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "circle": MSO_SHAPE.OVAL,
        "chevron": MSO_SHAPE.CHEVRON,
        "right_arrow": MSO_SHAPE.RIGHT_ARROW,
        "left_arrow": MSO_SHAPE.LEFT_ARROW,
        "pentagon": MSO_SHAPE.PENTAGON,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    }

    mso_shape = shape_map.get(subtype, MSO_SHAPE.RECTANGLE)
    shape = slide.shapes.add_shape(
        mso_shape, Inches(spec["x"]), Inches(spec["y"]),
        Inches(spec["w"]), Inches(spec["h"])
    )

    # Fill
    if "fill" in spec and spec["fill"]:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(spec["fill"])

    # Border
    border = spec.get("border")
    if border:
        shape.line.color.rgb = hex_to_rgb(border["color"])
        shape.line.width = Pt(border.get("width_pt", 1))
    else:
        shape.line.fill.background()

    # Corner radius (for rounded rectangles)
    if subtype == "rounded_rectangle" and "corner_radius_inches" in spec:
        radius = int(spec["corner_radius_inches"] * 914400)
        try:
            spPr = shape._element.find(qn('a:spPr')) or shape._element.find(qn('p:spPr'))
            if spPr is not None:
                prstGeom = spPr.find(qn('a:prstGeom'))
                if prstGeom is not None:
                    avLst = prstGeom.find(qn('a:avLst'))
                    if avLst is not None:
                        for gd in avLst.findall(qn('a:gd')):
                            if gd.get('name') == 'adj':
                                gd.set('fmla', f'val {radius}')
        except Exception:
            pass

    # Shadow
    shadow = spec.get("shadow")
    if shadow:
        try:
            spPr = shape._element.find(qn('a:spPr')) or shape._element.find(qn('p:spPr'))
            if spPr is not None:
                effectLst = spPr.find(qn('a:effectLst'))
                if effectLst is None:
                    effectLst = spPr.makeelement(qn('a:effectLst'), {})
                    spPr.append(effectLst)

                blur_emu = int(shadow.get("blur_pt", 4) * 12700)
                dist_emu = int(shadow.get("distance_pt", 2) * 12700)
                dir_val = shadow.get("direction_deg", 135) * 60000  # PPT uses 1/60000 deg
                shadow_hex = shadow.get("color", "#000000").lstrip("#")

                outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
                    'blurRad': str(blur_emu),
                    'dist': str(dist_emu),
                    'dir': str(dir_val),
                })
                srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': shadow_hex[:6]})
                alpha = outerShdw.makeelement(qn('a:alpha'), {'val': str(shadow.get("alpha_pct", 40000))})
                srgbClr.append(alpha)
                outerShdw.append(srgbClr)
                effectLst.append(outerShdw)
        except Exception:
            pass

    # Rotation
    if "rotation_deg" in spec:
        shape.rotation = spec["rotation_deg"]

    return shape


def add_line(slide, spec):
    """Add a line connector."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(spec["x1"]), Inches(spec["y1"]),
        Inches(spec["x2"]), Inches(spec["y2"])
    )
    connector.line.color.rgb = hex_to_rgb(spec.get("color", "#000000"))
    connector.line.width = Pt(spec.get("width_pt", 1))

    # Dash style if specified
    if spec.get("dash"):
        from pptx.oxml.ns import qn
        dash_map = {
            "dash": "dash",
            "dot": "dot",
            "dash_dot": "dashDot",
        }
        prstDash_str = dash_map.get(spec["dash"], "solid")
        try:
            ln = connector._element.find(qn('a:ln')) or connector._element.find(qn('p:ln'))
            if ln is not None:
                ln.set('prstDash', prstDash_str)
        except Exception:
            pass

    return connector


def add_simple_table(slide, spec):
    """Add a table to the slide."""
    rows = spec["rows"]
    cols = spec["cols"]

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(spec["x"]), Inches(spec["y"]),
        Inches(spec["w"]), Inches(spec["h"])
    )
    table = table_shape.table

    # Set column widths evenly
    col_width = Inches(spec["w"] / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    # Set row heights evenly
    row_height = Inches(spec["h"] / rows)
    for r in range(rows):
        table.rows[r].height = row_height

    # Default cell formatting
    default_font = spec.get("font", "PingFang SC")
    default_size = spec.get("size_pt", 10)
    default_color = spec.get("color", "#000000")
    border_color = spec.get("border_color")
    border_width = Pt(spec.get("border_width_pt", 0.5))

    # Apply default formatting to all cells first
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
            # Clear default text
            cell.text_frame.paragraphs[0].clear()
            # Set border
            if border_color:
                _set_cell_borders(cell, border_color, border_width)

    # Apply cell-specific formatting
    for cell_spec in spec.get("cells", []):
        r = cell_spec["row"]
        c = cell_spec["col"]
        if r >= rows or c >= cols:
            continue

        cell = table.cell(r, c)
        p = cell.text_frame.paragraphs[0]

        if "alignment" in cell_spec:
            p.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(cell_spec["alignment"], PP_ALIGN.CENTER)

        run = p.add_run()
        run.text = cell_spec.get("text", "")

        font_name = cell_spec.get("font", default_font)
        run.font.name = font_name
        try:
            rPr = run._r.get_or_add_rPr()
            ea = rPr.makeelement(qn('a:ea'), {})
            ea.set('typeface', font_name)
            rPr.append(ea)
        except Exception:
            pass

        run.font.size = Pt(cell_spec.get("size_pt", default_size))
        run.font.bold = cell_spec.get("bold", False)
        run.font.color.rgb = hex_to_rgb(cell_spec.get("color", default_color))

        # Cell fill
        if "fill" in cell_spec:
            try:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': cell_spec["fill"].lstrip("#")})
                solidFill.append(srgbClr)
                tcPr.append(solidFill)
            except Exception:
                pass

        # Merge cells
        if "merge" in cell_spec:
            merge_spec = cell_spec["merge"]
            other_cell = table.cell(merge_spec["row"], merge_spec["col"])
            cell.merge(other_cell)

    return table_shape


def _set_cell_borders(cell, color_hex, width):
    """Set borders on a table cell."""
    color_hex = color_hex.lstrip("#")
    try:
        tcPr = cell._tc.get_or_add_tcPr()
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = tcPr.makeelement(qn(border_name), {'w': str(int(width))})
            solidFill = ln.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)
    except Exception:
        pass


def build_slide(spec, output_path, template_path=None):
    """Build a single-slide PPTX from a JSON specification."""
    slide_spec = spec.get("slide", spec)

    # Create presentation
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Set slide dimensions
    if "width_inches" in slide_spec:
        prs.slide_width = Inches(slide_spec["width_inches"])
    if "height_inches" in slide_spec:
        prs.slide_height = Inches(slide_spec["height_inches"])

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Set background
    if "background" in slide_spec:
        set_slide_background(slide, slide_spec["background"])

    # Sort elements by layer — check both inside slide_spec and at root level
    elements = slide_spec.get("elements", []) or spec.get("elements", [])
    elements = sorted(elements, key=lambda e: e.get("layer", 0))

    # Build each element
    for elem in elements:
        etype = elem.get("type")
        try:
            if etype == "text_box":
                add_text_box(slide, elem)
            elif etype == "image":
                add_image(slide, elem)
            elif etype == "shape":
                add_shape(slide, elem)
            elif etype == "line":
                add_line(slide, elem)
            elif etype == "simple_table":
                add_simple_table(slide, elem)
            else:
                print(f"WARNING: Unknown element type: {etype}", file=sys.stderr)
        except Exception as e:
            print(f"ERROR building element {etype}: {e}", file=sys.stderr)
            import traceback
            traceback.print_exc()

    # Save
    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"PPTX saved to: {out_path.resolve()}")


def main():
    parser = argparse.ArgumentParser(description="Build a single-slide PPTX from JSON spec")
    parser.add_argument("--spec", required=True, help="Path to JSON specification file")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    parser.add_argument("--template", help="Optional template PPTX to use as base")
    args = parser.parse_args()

    if not Path(args.spec).exists():
        print(f"ERROR: Spec file not found: {args.spec}", file=sys.stderr)
        sys.exit(1)

    with open(args.spec, "r", encoding="utf-8") as f:
        spec = json.load(f)

    build_slide(spec, args.output, args.template)


if __name__ == "__main__":
    main()
