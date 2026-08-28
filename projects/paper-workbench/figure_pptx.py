# -*- coding: utf-8 -*-
"""figure_pptx.py — PPT 绘图路由（工作台固化工具）。

流程：python-pptx 程序化绘制示意图/流程图 → PowerPoint COM 导出高清 PNG。
配色：ggsci NPG（Nature 出版集团）科研配色，色盲友好。

用法：
    python figure_pptx.py fig2 <out.png>    # 放大过滤流程图
    python figure_pptx.py fig3 <out.png>    # 育种技术路线图
"""
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- ggsci NPG 配色（Nature 出版集团） ----------
NPG = {
    "red": "E64B35", "cyan": "4DBBD5", "green": "00A087", "navy": "3C5488",
    "salmon": "F39B7F", "greyblue": "8491B4", "lightteal": "91D1C2",
    "crimson": "DC0000", "brown": "7E6148", "khaki": "B09C85",
}
LIGHT = {"navy": "E8EDF5", "cyan": "E6F4F7", "green": "E4F3F0",
         "red": "FBEAE7", "grey": "F2F3F5"}
FONT = "Arial"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs, prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_box(slide, x, y, w, h, text, fill=None, line=None, font_color="222222",
            size=14, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.5,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = 0.12
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(font_color)
    return sp


def add_arrow(slide, x1, y1, x2, y2, color="666666", width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = RGBColor.from_string(color)
    conn.line.width = Pt(width)
    # 箭头（COM 级 XML）
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


# ---------- Figure 2: 放大过滤流程 ----------
def build_fig2(out_pptx):
    prs, slide = new_prs()
    add_box(slide, 0.45, 0.35, 12.4, 0.75,
            "From flask-level selection to industrial release: breeding target traits must pass the scale-up filter",
            fill=NPG["navy"], font_color="FFFFFF", size=16, bold=True)
    traits = [("Biomass\nproductivity", NPG["cyan"]),
              ("Marker\nmetabolites", NPG["green"]),
              ("Morphology &\nrheology", NPG["salmon"]),
              ("Genetic\nstability", NPG["greyblue"])]
    ty0 = 1.75
    for i, (t, c) in enumerate(traits):
        y = ty0 + i * 1.32
        add_box(slide, 0.85, y, 2.55, 1.06, t, fill=LIGHT["cyan"], line=c,
                size=13.5, bold=True, font_color="333333", line_w=2.0)
        add_arrow(slide, 3.45, y + 0.53, 5.15, 4.05, color="999999", width=1.6)
    add_box(slide, 5.2, 2.75, 3.1, 2.6,
            "SCALE-UP FILTER\n\nO\u2082 transfer\nShear tolerance\nPellet morphology\nBioreactor consistency",
            fill=NPG["navy"], font_color="FFFFFF", size=13.5, bold=True)
    add_arrow(slide, 8.35, 4.05, 9.95, 4.05, color=NPG["green"], width=2.6)
    add_box(slide, 10.0, 3.15, 2.7, 1.85,
            "RELEASED\ncandidate strain\n\nauthenticated &\nstable",
            fill=NPG["green"], font_color="FFFFFF", size=13.5, bold=True)
    add_box(slide, 0.85, 6.6, 11.85, 0.6,
            "Strains selected only under flask conditions may fail at production scale; scale-representative evaluation is a breeding-stage requirement.",
            fill=None, line=None, size=11.5, font_color="666666", align=PP_ALIGN.LEFT)
    prs.save(out_pptx)


# ---------- Figure 3: 育种技术路线图 ----------
def build_fig3(out_pptx):
    prs, slide = new_prs()
    add_box(slide, 0.6, 0.5, 12.15, 0.9,
            "INDUSTRIAL PULL:  biomass productivity   •   marker metabolites   •   fermentation robustness   •   genetic stability",
            fill=NPG["navy"], font_color="FFFFFF", size=14.5, bold=True)
    stages = [("Natural\nselection", NPG["greyblue"], LIGHT["grey"],
               "isolation & screening\nunindexed-heavy evidence"),
              ("Mutation\nbreeding", NPG["greyblue"], LIGHT["grey"],
               "UV / chemical / ARTP\nmostly unindexed in this species"),
              ("Protoplast\nfusion", NPG["red"], LIGHT["red"],
               "EVIDENCE GAP\nno indexed study in\nthis species to date"),
              ("Genome-informed\ndesign", NPG["green"], LIGHT["green"],
               "inflection: ATMT 2024\nmarkers & candidate genes")]
    sx0, sw, gap = 0.6, 2.82, 0.28
    for i, (name, c, lf, note) in enumerate(stages):
        x = sx0 + i * (sw + gap)
        add_box(slide, x, 2.05, sw, 1.35, name, fill=c, font_color="FFFFFF",
                size=16, bold=True)
        add_box(slide, x, 3.62, sw, 1.15, note, fill=lf, line=c, size=11,
                font_color="444444", line_w=1.75,
                bold=(i == 2))
        if i < 3:
            add_arrow(slide, x + sw + 0.02, 2.72, x + sw + gap - 0.02, 2.72,
                      color="888888", width=2.2)
    add_box(slide, 0.6, 5.45, 12.15, 0.9,
            "ENABLING BASE:  genome & transcriptome resources   •   multi-locus authentication   •   stability / cell-bank standards",
            fill=NPG["greyblue"], font_color="FFFFFF", size=14.5, bold=True)
    add_box(slide, 0.6, 6.65, 12.15, 0.6,
            "Priority: close the protoplast-fusion gap first; layer molecular tools onto fusion and mutation pipelines as they mature.",
            fill=None, line=None, size=11.5, font_color="666666", align=PP_ALIGN.LEFT)
    prs.save(out_pptx)


# ---------- 通用示意图构建器（figure_router 调用） ----------
def build_spec(spec, out_pptx):
    """通用示意图：spec = {
      "title": 顶部标题条(可选),
      "boxes": [{"x","y","w","h","text","fill"?,"line"?,"color"?,"size"?,"bold"?}],  # 单位: 英寸
      "arrows": [{"x1","y1","x2","y2","color"?}],
      "texts": [{"x","y","w","h","text","size"?,"color"?}],
      "footer": 底部说明(可选)
    } 坐标基于 13.333 x 7.5 英寸画布。"""
    prs, slide = new_prs()
    if spec.get("title"):
        add_box(slide, 0.45, 0.35, 12.4, 0.75, spec["title"],
                fill=NPG["navy"], font_color="FFFFFF", size=15, bold=True)
    for b in spec.get("boxes", []):
        add_box(slide, float(b["x"]), float(b["y"]), float(b["w"]), float(b["h"]),
                b.get("text", ""), fill=b.get("fill"), line=b.get("line"),
                font_color=b.get("color", "FFFFFF" if b.get("fill") else "333333"),
                size=float(b.get("size", 13)), bold=bool(b.get("bold", False)))
    for a in spec.get("arrows", []):
        add_arrow(slide, float(a["x1"]), float(a["y1"]), float(a["x2"]), float(a["y2"]),
                  color=a.get("color", "888888"), width=float(a.get("width", 2.0)))
    for tx in spec.get("texts", []):
        add_box(slide, float(tx["x"]), float(tx["y"]), float(tx.get("w", 6)), float(tx.get("h", 0.6)),
                tx.get("text", ""), fill=None, line=None,
                font_color=tx.get("color", "666666"), size=float(tx.get("size", 11.5)),
                align=PP_ALIGN.LEFT)
    if spec.get("footer"):
        add_box(slide, 0.6, 6.7, 12.15, 0.55, spec["footer"], fill=None, line=None,
                size=11, font_color="666666", align=PP_ALIGN.LEFT)
    prs.save(out_pptx)


# ---------- PowerPoint COM 导出 PNG ----------
def export_png(pptx_path, png_path, width_px=4000):
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        app.Visible = 1
    except Exception:
        pass
    pres = None
    try:
        try:
            pres = app.Presentations.Open(str(pptx_path), ReadOnly=True, Untitled=False, WithWindow=False)
        except Exception:
            pres = app.Presentations.Open(str(pptx_path), ReadOnly=True, Untitled=False, WithWindow=True)
        h_px = int(width_px * 7.5 / 13.333)
        pres.Slides(1).Export(str(png_path), "PNG", width_px, h_px)
    finally:
        if pres is not None:
            pres.Close()
        app.Quit()
    return png_path


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    kind, out = sys.argv[1], sys.argv[2]
    tmp_pptx = Path(out).with_suffix(".pptx")
    if kind == "fig2":
        build_fig2(str(tmp_pptx))
    elif kind == "fig3":
        build_fig3(str(tmp_pptx))
    elif kind == "generic":
        import json as _json
        spec = _json.loads(Path(sys.argv[3]).read_text(encoding="utf-8")) if len(sys.argv) > 3 else {}
        build_spec(spec, str(tmp_pptx))
    else:
        print("未知类型:", kind); sys.exit(1)
    export_png(tmp_pptx, out, width_px=int(sys.argv[3]) if len(sys.argv) > 3 else 4000)
    print("OK:", out, os.path.getsize(out), "bytes; 源 PPTX:", tmp_pptx)


if __name__ == "__main__":
    main()
